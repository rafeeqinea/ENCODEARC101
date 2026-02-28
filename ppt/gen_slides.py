from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    def add_slide(title_text, subtitle_text, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title_text
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = subtitle_text
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'ArcTreasury: AI-Powered Autonomous Treasury Management'
    subtitle.text = 'Encode x Arc Enterprise & DeFi Hackathon'

    # Slide 2: The Vision
    add_slide('1. The Vision', 'Autonomous Liquidity Excellence.', [
        'ArcTreasury is an AI-driven agent that manages corporate stablecoin reserves on the Arc Network.',
        'Ensuring optimal yield, seamless FX hedging, and automated obligation fulfillment.'
    ])

    # Slide 3: The Problem
    add_slide('2. The Problem', 'Traditional treasury management is inefficient.', [
        'Fragmentation: Liquidity split across USDC and EURC.',
        'Idle Capital: Assets sitting in wallets earning 0% yield.',
        'Manual Burden: Human teams manually swapping and paying invoices.'
    ])

    # Slide 4: Our Solution: The AI Agent
    add_slide('3. Our Solution: The AI Agent', 'A 24/7 autonomous loop that handles the heavy lifting.', [
        'Monitors: Real-time FX rates (Circle StableFX) and Yield (USYC).',
        'Predicts: ML-based FX forecasting and risk assessment.',
        'Executes: On-chain swaps and yield deposits on the Arc Testnet.'
    ])

    # Slide 5: Technical Architecture
    add_slide('4. Technical Architecture', 'Built on cutting-edge stacks.', [
        'Frontend: React + Framer Motion (Real-time Liquid Glass Dashboard).',
        'AI Backend: Python FastAPI Agent + predictive ML models.',
        'Smart Contracts: Solidity Vaults on Arc Testnet for asset custody.',
        'Integrations: Circle StableFX, Hashnote USYC, Stork Oracle.'
    ])

    # Slide 6: Key Demo Features
    add_slide('5. Key Demo Features', 'Real tools, real results.', [
        'Live FX Monitor: Execute trades with real-time quotes.',
        'Automated Obligations: Agent pre-funds EURC ahead of pay dates.',
        'Yield Optimization: Idle USDC is automatically parked in USYC for 4.5% APY.'
    ])

    # Slide 7: The Roadmap
    add_slide('6. The Roadmap', 'The future of autonomous finance.', [
        'Multi-chain Expansion: Using Circle CCTP for cross-chain rebalancing.',
        'Advanced ML: Reinforcement learning for complex hedging strategies.',
        'Institutional Controls: Multi-sig approval hooks for high-value agent decisions.'
    ])

    output_path = 'ArcTreasury_Pitch.pptx'
    prs.save(output_path)
    print(f'PPTX saved to {os.path.abspath(output_path)}')

if __name__ == "__main__":
    create_presentation()
